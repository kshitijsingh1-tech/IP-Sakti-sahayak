"""
IP-SAKTI Document Processor
-----------------------------
Handles PDF, PPT, images and plain text files uploaded by the user.
Extracts raw text/content and prepares a structured context string
that the 4-Agent RAG pipeline can consume as additional evidence.
"""
from __future__ import annotations

import io
import os
import re
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional heavy dependencies — gracefully degrade if not installed
# ---------------------------------------------------------------------------
try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from PIL import Image as PILImage
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


# ---------------------------------------------------------------------------
# Public Interface
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".pdf", ".ppt", ".pptx",
    ".doc", ".docx",
    ".txt", ".csv",
    ".jpg", ".jpeg", ".png", ".webp", ".gif"
}

MAX_EXTRACTED_CHARS = 8_000  # Characters sent to LLM context per document


def extract_text(filename: str, content: bytes) -> str:
    """
    Dispatch to the correct extractor based on file extension.
    Returns a cleaned text string ready for RAG pipeline injection.
    Raises ValueError for unsupported file types.
    """
    ext = Path(filename).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Allowed: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    try:
        if ext == ".pdf":
            text = _extract_pdf(content)
        elif ext in (".ppt", ".pptx"):
            text = _extract_pptx(content)
        elif ext in (".doc", ".docx"):
            text = _extract_docx(content)
        elif ext == ".csv":
            text = _extract_csv(content)
        elif ext == ".txt":
            text = _extract_txt(content)
        elif ext in (".jpg", ".jpeg", ".png", ".webp", ".gif"):
            text = _extract_image(content)
        else:
            text = ""
    except Exception as exc:
        logger.warning("Extraction failed for %s: %s", filename, exc)
        text = f"[Extraction failed for {filename}: {exc}]"

    cleaned = _clean(text)
    return cleaned[:MAX_EXTRACTED_CHARS]


def build_document_context(filename: str, text: str) -> str:
    """
    Wraps extracted text in a structured prompt block for the RAG agents.
    """
    return (
        f"=== ATTACHED DOCUMENT: {filename} ===\n"
        f"{text}\n"
        f"=== END OF DOCUMENT ===\n"
    )


# ---------------------------------------------------------------------------
# Individual Extractors
# ---------------------------------------------------------------------------

def _extract_pdf(content: bytes) -> str:
    if not PYPDF_AVAILABLE:
        return "[pypdf not installed — run: pip install pypdf]"
    reader = pypdf.PdfReader(io.BytesIO(content))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


def _extract_pptx(content: bytes) -> str:
    if not PPTX_AVAILABLE:
        return "[python-pptx not installed — run: pip install python-pptx]"
    prs = PptxPresentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_docx(content: bytes) -> str:
    if not DOCX_AVAILABLE:
        return "[python-docx not installed — run: pip install python-docx]"
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_csv(content: bytes) -> str:
    if not PANDAS_AVAILABLE:
        # Plain text fallback
        return content.decode("utf-8", errors="replace")
    try:
        df = pd.read_csv(io.BytesIO(content))
        return df.to_string(index=False, max_rows=200)
    except Exception:
        return content.decode("utf-8", errors="replace")


def _extract_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _extract_image(content: bytes) -> str:
    """
    OCR via Tesseract if available; otherwise returns a placeholder
    instructing the LLM that an image was attached.
    """
    if not OCR_AVAILABLE:
        return (
            "[Image attached — OCR not available. "
            "Install: pip install pillow pytesseract and configure Tesseract. "
            "The image may contain formulation labels, patent drawings, or clinical data.]"
        )
    img = PILImage.open(io.BytesIO(content))
    return pytesseract.image_to_string(img)


# ---------------------------------------------------------------------------
# Text Cleaning
# ---------------------------------------------------------------------------

def _clean(text: str) -> str:
    """Remove excessive whitespace and non-printable characters."""
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    # Strip non-printable chars (keep Unicode letters/symbols)
    text = "".join(ch for ch in text if ch.isprintable() or ch == '\n')
    return text.strip()
